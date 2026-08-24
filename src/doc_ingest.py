#!/usr/bin/env python3
"""文档识别摄入（2026-08-22 用户方案：夜班识别微信/邮件收到的文档文件）。

扫描 exchange/inbox/** 新文档（pdf/docx/doc/xlsx/pptx/txt/md），提取文本 →
L0(source=doc:file) → 文本自动进 L2 索引。扫描件/提取失败 → Qwen3-VL-4B OCR 兜底。

提取器：
  docx/doc → macOS textutil（系统自带）
  pdf     → pypdf 提取；空文本（扫描件）→ sips 转 png → VL-4B OCR
  xlsx    → openpyxl（read_only，截断）
  pptx    → python-pptx（未装则跳过并记录）
  txt/md  → 直接读

运行环境：doc-parse venv（pypdf/openpyxl/python-docx）。沙盒：AIAGENT_SANDBOX 跟随。
用法：
  python3 doc_ingest.py          # 扫 + 提取 + 落 L0（夜班 seg1 后调用）
"""
from __future__ import annotations

import hashlib
import json
import os
import subprocess
import sys
import time

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from l0_ingest import L0Writer  # noqa: E402

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
EXCHANGE = os.path.join(REPO, "exchange")
L0_ROOT = os.path.join(REPO, "memory", "L0_raw")
STATE_PATH = os.path.join(REPO, "memory", "L1_working", "doc_ingest_state.json")

SANDBOX = os.environ.get("AIAGENT_SANDBOX", "")
if SANDBOX:
    EXCHANGE = os.path.join(SANDBOX, "exchange")
    L0_ROOT = os.path.join(SANDBOX, "memory", "L0_raw")
    STATE_PATH = os.path.join(SANDBOX, "memory", "L1_working", "doc_ingest_state.json")

DOC_EXTS = (".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".txt", ".md")
MAX_CHARS = 6000  # 单文件提取上限（防大文件爆炸）


def _extract_text(p: str) -> str:
    """按扩展名提取文本。返回截断文本（失败返回空串）。"""
    ext = os.path.splitext(p)[1].lower()
    try:
        if ext in (".txt", ".md"):
            with open(p, encoding="utf-8", errors="replace") as f:
                return f.read()[:MAX_CHARS]
        if ext in (".docx", ".doc"):
            r = subprocess.run(["textutil", "-convert", "txt", "-stdout", p],
                               capture_output=True, text=True, timeout=60)
            return r.stdout[:MAX_CHARS] if r.returncode == 0 else ""
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(p)
            text = "\n".join((pg.extract_text() or "") for pg in reader.pages[:20])
            return text[:MAX_CHARS]
        if ext in (".xlsx", ".xls"):
            from openpyxl import load_workbook
            wb = load_workbook(p, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets[:5]:
                rows = []
                for row in ws.iter_rows(max_row=200, values_only=True):
                    vals = [str(v) for v in row if v is not None]
                    if vals:
                        rows.append(" | ".join(vals))
                if rows:
                    parts.append(f"[sheet:{ws.title}]\n" + "\n".join(rows))
                if sum(len(x) for x in parts) > MAX_CHARS:
                    break
            wb.close()
            return "\n".join(parts)[:MAX_CHARS]
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(p)
            parts = []
            for slide in prs.slides[:30]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)[:MAX_CHARS]
    except Exception as e:
        return f"(提取失败 {type(e).__name__})"
    return ""


def _ocr_fallback(p: str) -> str:
    """PDF 无文本（扫描件）→ sips 转 png → Qwen3-VL-4B OCR。返回识别文本。"""
    try:
        png = p + ".ocr.png"
        subprocess.run(["sips", "-s", "format", "png", p, "--out", png],
                       capture_output=True, timeout=60)
        if not os.path.exists(png):
            return ""
        import base64, urllib.request
        b64 = base64.b64encode(open(png, "rb").read()).decode()
        payload = json.dumps({
            "model": "qwen3-vl",
            "messages": [{"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                {"type": "text", "text": "提取这个文档页的所有文字，原样输出。"}]}],
            "max_tokens": 800, "temperature": 0.0,
        }).encode()
        req = urllib.request.Request("http://127.0.0.1:8081/v1/chat/completions",
                                     data=payload, headers={"Content-Type": "application/json"},
                                     method="POST")
        resp = json.loads(urllib.request.urlopen(req, timeout=300).read())
        os.remove(png)
        return (resp["choices"][0]["message"].get("content") or "")[:MAX_CHARS]
    except Exception:
        return ""


def _ensure_vl() -> bool:
    try:
        import urllib.request
        urllib.request.urlopen(urllib.request.Request("http://127.0.0.1:8081/v1/models"), timeout=5)
        return True
    except Exception:
        pass
    script = os.path.join(REPO, "src", "serve_vision.sh")
    if not os.path.exists(script):
        return False
    subprocess.Popen(["bash", script, "8081"], stdout=subprocess.DEVNULL,
                     stderr=subprocess.DEVNULL, cwd=REPO, start_new_session=True)
    for _ in range(60):
        time.sleep(3)
        try:
            import urllib.request
            urllib.request.urlopen(urllib.request.Request("http://127.0.0.1:8081/v1/models"), timeout=3)
            return True
        except Exception:
            pass
    return False


def run() -> int:
    w = L0Writer(L0_ROOT)
    state = {"files": {}}
    if os.path.exists(STATE_PATH):
        try:
            state = json.load(open(STATE_PATH, encoding="utf-8"))
        except Exception:
            pass
    n = 0
    for dirpath, _dirs, files in os.walk(os.path.join(EXCHANGE, "inbox")):
        # 跳过邮件流程产物目录（email_deep 的 邮箱深度-*.md / 摘要已在 email 源入 L0，
        # 再摄入 = 重复。2026-08-23 修）
        if os.path.basename(dirpath) == "email" or "/email/" in dirpath.replace(EXCHANGE, ""):
            continue
        for fn in sorted(files):
            if not fn.lower().endswith(DOC_EXTS) or fn.startswith("."):
                continue
            p = os.path.join(dirpath, fn)
            rel = os.path.relpath(p, REPO)
            try:
                sig = hashlib.md5(open(p, "rb").read()).hexdigest()
            except OSError:
                continue
            if state["files"].get(rel) == sig:
                continue
            text = _extract_text(p)
            if text.startswith("(提取失败") or not text.strip():
                # 扫描件/失败 → VL-4B OCR（仅 pdf 转图可用时）
                if _ensure_vl():
                    ocr = _ocr_fallback(p) if p.lower().endswith(".pdf") else ""
                    text = text if text.strip() else ocr
            if not text.strip():
                print(f"[doc] 无文本 {rel}，跳过")
                state["files"][rel] = sig
                continue
            w.append("doc:file", {"path": rel, "text": text[:MAX_CHARS]},
                     sensitive=False, meta={"ingest": "doc_ingest", "ext": os.path.splitext(fn)[1].lower()})
            state["files"][rel] = sig
            n += 1
            print(f"[doc] {rel} → L0（{len(text)} chars）: {text[:50]}…")
    with open(STATE_PATH, "w", encoding="utf-8") as f:
        json.dump(state, f, ensure_ascii=False, indent=1)
    print(f"[doc] 完成 {n} 个文档")
    return n


if __name__ == "__main__":
    run()
